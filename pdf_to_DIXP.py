import tkinter as tk
from tkinter import filedialog, messagebox
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from pdf2docx import Converter
import pandas as pd
import tabula
import pytesseract
import os

def trim_image(image):
    """Trim the white space from the image."""
    bbox = image.getbbox()
    if bbox:
        return image.crop(bbox)
    return image

def convert_pdf_to_images(pdf_path):
    try:
        # Convert PDF to images (list of images, one per page)
        images = convert_from_path(pdf_path)

        # Save images to the same directory as the PDF
        output_folder = os.path.splitext(pdf_path)[0] + "_images"
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        for i, image in enumerate(images):
            # Trim the image to remove white space
            trimmed_image = trim_image(image)
            image_path = os.path.join(output_folder, f"page_{i + 1}.png")
            trimmed_image.save(image_path, 'PNG')

        messagebox.showinfo("Success", f"PDF successfully converted to images in: {output_folder}")
    except Exception as e:
        messagebox.showerror("Error", f"Failed to convert PDF to images: {e}")

def convert_pdf_to_docx(pdf_path):
    try:
        if messagebox.askyesno("Scanned PDF", "Is this a scanned PDF?"):
            # Perform OCR
            images = convert_from_path(pdf_path)
            ocr_text = ""
            for img in images:
                text = pytesseract.image_to_string(img)
                ocr_text += text + "\n"
            
            docx_path = os.path.splitext(pdf_path)[0] + "_scanned.docx"
            with open(docx_path, "w") as docx_file:
                docx_file.write(ocr_text)
            messagebox.showinfo("Success", f"Scanned PDF successfully converted to DOCX: {docx_path}")
        else:
            # Convert PDF to DOCX directly
            docx_path = os.path.splitext(pdf_path)[0] + ".docx"
            cv = Converter(pdf_path)
            cv.convert(docx_path, start=0, end=None)
            cv.close()
            messagebox.showinfo("Success", f"PDF successfully converted to DOCX: {docx_path}")
    except Exception as e:
        messagebox.showerror("Error", f"Failed to convert PDF to DOCX: {e}")

def convert_pdf_to_excel(pdf_path):
    try:
        excel_path = os.path.splitext(pdf_path)[0] + ".xlsx"
        tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
        
        if not tables:
            raise Exception("No tables found in the PDF.")
        
        with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
            for i, table in enumerate(tables):
                table.to_excel(writer, sheet_name=f"Table_{i + 1}", index=False)
        
        messagebox.showinfo("Success", f"PDF successfully converted to Excel: {excel_path}")
    except Exception as e:
        messagebox.showerror("Error", f"Failed to convert PDF to Excel: {e}")

def convert_pdf_to_pptx(pdf_path):
    try:
        images = convert_from_path(pdf_path)
        ppt = Presentation()

        for image in images:
            trimmed_image = trim_image(image)
            img_path = os.path.splitext(pdf_path)[0] + "_temp.png"
            trimmed_image.save(img_path, "PNG")

            slide_layout = ppt.slide_layouts[5]
            slide = ppt.slides.add_slide(slide_layout)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                      width=ppt.slide_width, height=ppt.slide_height)
            os.remove(img_path)

        ppt_path = os.path.splitext(pdf_path)[0] + ".pptx"
        ppt.save(ppt_path)
        messagebox.showinfo("Success", f"PDF successfully converted to PPTX: {ppt_path}")
    except Exception as e:
        messagebox.showerror("Error", f"Failed to convert PDF to PPTX: {e}")

def select_file():
    pdf_path = filedialog.askopenfilename(
        title="Select PDF file",
        filetypes=[("PDF files", "*.pdf")]
    )
    if pdf_path:
        choice = conversion_choice.get()
        if choice == "1":
            convert_pdf_to_images(pdf_path)
        elif choice == "2":
            convert_pdf_to_docx(pdf_path)
        elif choice == "3":
            convert_pdf_to_excel(pdf_path)
        elif choice == "4":
            convert_pdf_to_pptx(pdf_path)

# Create Tkinter window
root = tk.Tk()
root.title("PDF Converter")

# Add a label and radio buttons for conversion choice
conversion_choice = tk.StringVar(value="1")  # Default selection
tk.Label(root, text="Select conversion type:").pack(pady=10)
tk.Radiobutton(root, text="1. PDF to Images", variable=conversion_choice, value="1").pack(anchor=tk.W)
tk.Radiobutton(root, text="2. PDF to DOCX", variable=conversion_choice, value="2").pack(anchor=tk.W)
tk.Radiobutton(root, text="3. PDF to Excel", variable=conversion_choice, value="3").pack(anchor=tk.W)
tk.Radiobutton(root, text="4. PDF to PPTX", variable=conversion_choice, value="4").pack(anchor=tk.W)

# Create and place a button to trigger file selection
select_button = tk.Button(root, text="Select PDF File", command=select_file)
select_button.pack(pady=20)

# Run the Tkinter event loop
root.mainloop()
